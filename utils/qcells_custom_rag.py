from llama_index.core.vector_stores import MetadataFilters, MetadataFilter,FilterOperator
from llama_index.core.response_synthesizers import ResponseMode
from llama_index.core.prompts import PromptTemplate
from llama_index.embeddings.azure_openai import AzureOpenAIEmbedding
from llama_index.core import SummaryIndex, get_response_synthesizer, VectorStoreIndex, ServiceContext, GPTVectorStoreIndex, Document
from pydantic import Field, BaseModel
from typing import List, Optional
from llama_index.core.tools.tool_spec.base import BaseToolSpec
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.tools import QueryEngineTool, ToolMetadata
from llama_index.core.query_engine import SubQuestionQueryEngine
from llama_index.core.schema import NodeWithScore, Node
from web_catch import * 
from llama_index.core.node_parser import CodeSplitter
from llama_index.llms.azure_openai import AzureOpenAI
import sys
sys.path.append("../utils/pytube")
from generate_summary import map_reduced_summary
import io
from pptx import Presentation
import fitz
from pytube import Channel,YouTube #키워드 -> url
from llama_index.core.tools import RetrieverTool
from llama_index.core.selectors import (
    PydanticMultiSelector,
    PydanticSingleSelector,
)
from llama_index.core.retrievers import RouterRetriever


llm =AzureOpenAI(
            model="gpt-35-turbo",
            deployment_name="qcell-gpt-model-rag",
            temperature = 0,
            api_key="c11ed4df2d35412b89a7b51a631bf0e4",
            azure_endpoint="https://rag-openai-qcells-east.openai.azure.com/",
            api_version="2024-02-15-preview",
        )
embedding = AzureOpenAIEmbedding(
        model="text-embedding-ada-002",
        deployment_name="qcell_embedding_model",
        api_key="c11ed4df2d35412b89a7b51a631bf0e4",
        azure_endpoint="https://rag-openai-qcells-east.openai.azure.com/",
     api_version="2023-07-01-preview")

def pptx_load_data(file):
    bytes_data = file.getvalue()
    byte_io = io.BytesIO(bytes_data)
    presentation = Presentation(byte_io)
    result = ""
    for i, slide in enumerate(presentation.slides):
        result += f"\n\nSlide #{i}: \n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                result += f"{shape.text}\n"                
    return result

def pdf_load_data(file):
    docs = []  
    file.seek(0,0)
    extra_info = {}
    doc = fitz.open(stream=file.read(), filetype="pdf")
    extra_info["total_pages"] = len(doc)
    extra_info['title'] = file.name
    all_text = '\n\n'.join([page.get_text()for page in doc])    
    return all_text

def get_youtube_metadata(url):
    yt = YouTube(url)
    video_info = {
        "url":url,
        "title": yt.title or "Unknown",
        "description": yt.description or "Unknown",
        "view_count": yt.views or 0,
        "thumbnail_url": yt.thumbnail_url or "Unknown",
        "publish_date": yt.publish_date.strftime("%Y-%m-%d %H:%M:%S")
        if yt.publish_date
        else "Unknown",
        "length": yt.length or 0,
        "author": yt.author or "Unknown",
        "embed_url" : yt.embed_url
    }
    return video_info

def get_timeline(seconds):
    m, s = divmod(seconds, 60)
    h, m = divmod(m, 60)
    return '{}:{}:{}'.format(int(h), format(int(m), '02'), format(int(s), '02'))

class Decide_to_search_external_web(BaseModel):
    """
        Verifying if the response to the question is accurate, and deciding whether to conduct an external search.
        return
            Succeed_answer (bool) : check wether answer is queried well
            Decide_web_search (bool): Assessing the request for the latest information in the question
            Searchable_query (str): Assessing the question in terms of easily searchable keywords on Google.(limited 4-words)
            Reason (bool): if query is clear question.
    """
    Succeed_answer: bool
    Decide_web_search: bool
    Searchable_query: str
    Reason: bool

mrs = map_reduced_summary(llm = llm, embedding = embedding)
class create_db_chat(BaseToolSpec):
    spec_functions = ['retriever_documents']
    
    def __init__(self, _docs, llm, embedding, service_context):
        self.llm = llm
        self.embedding = embedding
        self.service_context = service_context
        self.splitter = SentenceSplitter(chunk_size=256,chunk_overlap=25)
        self.splitter2 = SentenceSplitter(chunk_size=512,chunk_overlap=50)
        self.splitter_code = CodeSplitter(language="python", chunk_lines=40, chunk_lines_overlap=15, max_chars=1500)        
        self.response_synthesizer = get_response_synthesizer(llm = self.llm, response_mode="tree_summarize", use_async=True)
        self.query_engine_tools = []
        self.summary = []

        for idx, doc in enumerate(_docs):
            idx = str(idx)
            if doc[0].metadata['title'].split('.')[-1] == 'py':
                self.index_node = GPTVectorStoreIndex.from_documents(doc, service_context = self.service_context, transformations= [self.splitter_code],show_progress=True)            
                self.queryengine = self.index_node.as_query_engine()                
                self.query_engine_tools.append(RetrieverTool(query_engine= self.queryengine, 
                                                             name="Python_script_" + idx,
                                                             description= "Please answer questions about the content of the this python script"))
                self.summary = ''
    
            elif doc[0].metadata['title'].split('.')[-1] == 'pdf':
                self.index_node = GPTVectorStoreIndex.from_documents(doc, service_context = self.service_context, transformations= [self.splitter],show_progress=True)            
                self.retriever_engine = self.index_node.as_retriever()
                self.query_engine_tools.append(RetrieverTool.from_defaults(retriever=self.retriever_engine,
                                                                            name="pdf" + idx,
                                                                            description="Please answer questions about the content of the {}".format(doc[0].metadata['title'])))
                                               
                self.summary.append(doc[0].metadata['title'] + '\n\n' + mrs.create_document_summary('\n'.join([i.text for i in doc])) + '\n\n\n')
                
            elif doc[0].metadata['title'].split('.')[-1] == 'pptx':
                self.index_node = GPTVectorStoreIndex.from_documents(doc, service_context = self.service_context, transformations= [self.splitter],show_progress=True)            
                self.retriever_engine = self.index_node.as_retriever()
                self.query_engine_tools.append(RetrieverTool.from_defaults(retriever=self.retriever_engine,
                                                                            name="pptx_" + idx,
                                                                            description="Please answer questions about the content of the {}".format(doc[0].metadata['title'])))
                self.summary.append(doc[0].metadata['title'] + '\n\n' + mrs.create_document_summary('\n'.join([i.text for i in doc])) + '\n\n\n')
                
            elif doc[0].metadata['resource'] =='web_page': 
                self.index_node = GPTVectorStoreIndex.from_documents(doc, service_context = self.service_context, transformations= [self.splitter],show_progress=True)            
                self.retriever_engine = self.index_node.as_retriever()
                self.query_engine_tools.append(RetrieverTool.from_defaults(retriever=self.retriever_engine,
                                                                            name="web_html_" + idx,
                                                                            description="Please answer questions about the content of the {}".format(doc[0].metadata['title'])))
                self.summary.append(doc[0].metadata['title'] + '\n\n' + mrs.create_document_summary('\n'.join([i.text for i in doc])) + '\n\n\n')                
            else:
                self.index_node = GPTVectorStoreIndex.from_documents(doc, service_context = self.service_context, transformations= [self.splitter],show_progress=True)            
                self.retriever_engine = self.index_node.as_retriever()
                self.query_engine_tools.append(RetrieverTool.from_defaults(retriever=self.retriever_engine,
                                                                            name="youtube_transcript_" + idx,
                                                                            description="Please answer questions about the content of the {}".format(doc[0].metadata['title'])))
                self.summary.append(doc[0].metadata['title'] + '\n\n' + mrs.create_document_summary('\n'.join([i.text for i in doc])) + '\n\n\n')            
            
    def retriever_documents(self, query:str):
        """
        Answer a query about extracted documents.
        Return retrieved texts.

        Args:
            query (str): question about extracted documents.
        """
        retriever = RouterRetriever(
            selector=PydanticMultiSelector.from_defaults(llm=self.llm),
            llm = self.llm, 
            retriever_tools=self.query_engine_tools)
        nodes = retriever.retrieve(query)
        return [i.text for i in nodes]
