from llama_index.core.vector_stores import MetadataFilters, MetadataFilter,FilterOperator
from llama_index.core.response_synthesizers import ResponseMode
from llama_index.core.prompts import PromptTemplate
from llama_index.embeddings.azure_openai import AzureOpenAIEmbedding
from llama_index.core import SummaryIndex, get_response_synthesizer, VectorStoreIndex, ServiceContext, GPTVectorStoreIndex, Document
from pydantic import Field, BaseModel
from typing import List, Optional
from llama_index.core.tools.tool_spec.base import BaseToolSpec
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.tools import QueryEngineTool, ToolMetadata
from llama_index.core.query_engine import SubQuestionQueryEngine
from llama_index.core.schema import NodeWithScore, Node
from web_catch import * 
from llama_index.core.node_parser import CodeSplitter
from llama_index.llms.azure_openai import AzureOpenAI
import sys
sys.path.append("../utils/pytube")
from generate_summary import map_reduced_summary
import io
from pptx import Presentation
import fitz
from pytube import Channel,YouTube #키워드 -> url
from llama_index.core.tools import RetrieverTool
from llama_index.core.llms import ChatMessage, MessageRole
from llama_index.core.selectors import (
    PydanticMultiSelector,
    PydanticSingleSelector,
)
from llama_index.core.retrievers import RouterRetriever
from llama_index.program.openai import OpenAIPydanticProgram
from llama_index.retrievers.bm25 import BM25Retriever
from llama_index.core.retrievers import BaseRetriever
from llama_index.agent.openai_legacy import FnRetrieverOpenAIAgent
from llama_index.core.objects import ObjectIndex,SimpleToolNodeMapping,ObjectRetriever
from llama_index.postprocessor.cohere_rerank import CohereRerank
from llama_index.core.query_engine import RetrieverQueryEngine
from llama_index.core import get_response_synthesizer
from docx import Document as py_Document

llm =AzureOpenAI(
            model="gpt-35-turbo",
            deployment_name="qcell-gpt-model-rag",
            temperature = 0,
            api_key="c11ed4df2d35412b89a7b51a631bf0e4",
            azure_endpoint="https://rag-openai-qcells-east.openai.azure.com/",
            api_version="2024-02-15-preview",
        )
embedding = AzureOpenAIEmbedding(
        model="text-embedding-ada-002",
        deployment_name="qcell_embedding_model",
        api_key="c11ed4df2d35412b89a7b51a631bf0e4",
        azure_endpoint="https://rag-openai-qcells-east.openai.azure.com/",
     api_version="2023-07-01-preview")


class GenResult(BaseModel):
    """
    Get text data to list
    return
        search_task (list): 
        analyze_task: (list):
    """
    search_task: Optional [list]
    analyze_task: Optional [list]
    
def generate_strategy(high_level_query, answers):
    gen_program = OpenAIPydanticProgram.from_defaults(
    output_cls=GenResult,
    llm=llm,
    verbose=True,
    prompt_template_str=(
            "you will be given decomposed tasks and catgorized.\n"
            "Please make a list by category\n"
            "{context}"
        ),
    )   
    template = (
        "You are an tech sensing assistant to make a decision and plan for researching\n"
        "Please decompse the given query as much as possible you can break down and make tasks list as many as possible to let people know about tech."
        "Please categorize the decomposed tasks into two categories: data search and analyze."
        "the each task is limited 5-words and change searchable word on google"
        "\n---------------------\n"
        "The high level goal: {high_level_query}\n"
        "\n---------------------\n"
        "The previous answers: {answers}\n"
    )
    qa_template = PromptTemplate(template)
    prompt = qa_template.format(high_level_query = high_level_query, answers = answers)
    prompt = [ChatMessage(role='user', content=prompt)]
    res_tasks = llm.chat(prompt)    
    lists = gen_program(query = high_level_query, context = res_tasks.message.content)        
    return lists.search_task, lists.analyze_task

def docx_load_data(file):
    print(file)
    bytes_data = file.getvalue()
    byte_io = io.BytesIO(bytes_data)
    
    doc = py_Document(byte_io)
    text_list = []
    for paragraph in doc.paragraphs:
        text_list.append(paragraph.text)
    text_list = '\n'.join(text_list)
    return text_list


def pptx_load_data(file):
    bytes_data = file.getvalue()
    byte_io = io.BytesIO(bytes_data)
    presentation = Presentation(byte_io)
    result = ""
    for i, slide in enumerate(presentation.slides):
        result += f"\n\nSlide #{i}: \n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                result += f"{shape.text}\n"                
    return result

def pdf_load_data(file):
    docs = []  
    file.seek(0,0)
    extra_info = {}
    doc = fitz.open(stream=file.read(), filetype="pdf")
    extra_info["total_pages"] = len(doc)
    extra_info['title'] = file.name
    all_text = '\n\n'.join([page.get_text()for page in doc])    
    return all_text

def get_youtube_metadata(url):
    yt = YouTube(url)
    video_info = {
        "url":url,
        "title": yt.title or "Unknown",
        "description": yt.description or "Unknown",
        "view_count": yt.views or 0,
        "thumbnail_url": yt.thumbnail_url or "Unknown",
        "publish_date": yt.publish_date.strftime("%Y-%m-%d %H:%M:%S")
        if yt.publish_date
        else "Unknown",
        "length": yt.length or 0,
        "author": yt.author or "Unknown",
        "embed_url" : yt.embed_url
    }
    return video_info

def get_timeline(seconds):
    m, s = divmod(seconds, 60)
    h, m = divmod(m, 60)
    return '{}:{}:{}'.format(int(h), format(int(m), '02'), format(int(s), '02'))

def translator(llm, other_lang): 
    template = (
        "You are an AI agent to translate other language into English."
        "Translate into English, if the given text is not written in English."
        "Return just tranlated word"
        "\n---------------------\n"
        "the given text: {query}"
    )
    qa_template = PromptTemplate(template)
    prompt = qa_template.format(query = other_lang)
    prompt = [ChatMessage(role='user', content=prompt)]
    translated = llm.chat(prompt)
    return translated.message.content
    
class HybridRetriever(BaseRetriever):
    def __init__(self, vector_retriever, bm25_retriever):
        self.vector_retriever = vector_retriever
        self.bm25_retriever = bm25_retriever
        super().__init__()

    def _retrieve(self, query, **kwargs):
        bm25_nodes = self.bm25_retriever.retrieve(query, **kwargs)
        vector_nodes = self.vector_retriever.retrieve(query, **kwargs)

        all_nodes = []
        node_ids = set()
        for n in bm25_nodes + vector_nodes:
            if n.node.node_id not in node_ids:
                all_nodes.append(n)
                node_ids.add(n.node.node_id)
        return all_nodes

mrs = map_reduced_summary(llm = llm, embedding = embedding)
class CustomRetriever(BaseRetriever):
    def __init__(self, vector_retriever, postprocessor=None):
        self._vector_retriever = vector_retriever
        self._postprocessor = postprocessor 
        super().__init__()

    def _retrieve(self, query_bundle):
        retrieved_nodes = self._vector_retriever.retrieve(query_bundle)
        filtered_nodes = self._postprocessor.postprocess_nodes(
            retrieved_nodes, query_bundle=query_bundle
        )
        return filtered_nodes

class CustomObjectRetriever(ObjectRetriever):
    def __init__(self, retriever, object_node_mapping, all_tools, llm=None):
        self._retriever = retriever
        self._object_node_mapping = object_node_mapping
        self._llm = llm 

    def retrieve(self, query_bundle):
        nodes = self._retriever.retrieve(query_bundle)
        tools = [self._object_node_mapping.from_node(n.node) for n in nodes]

        sub_question_engine = SubQuestionQueryEngine.from_defaults(
            query_engine_tools=tools, llm=self._llm
        )
        sub_question_description = f"""\
Useful for any queries that involve comparing multiple documents. ALWAYS use this tool for comparison queries - make sure to call this \
tool with the original query. Do NOT use the other tools for any queries involving multiple documents.
"""
        sub_question_tool = QueryEngineTool(
            query_engine=sub_question_engine,
            metadata=ToolMetadata(
                name="compare_tool", description=sub_question_description
            ),
        )
        return tools + [sub_question_tool]
        
class create_db_chat(BaseToolSpec):
    spec_functions = ['hybrid_retriever_documents']
    
    def __init__(self, _docs, llm, embedding, service_context):
        self.llm = llm
        self._docs = _docs
        self.embedding = embedding
        self.service_context = service_context
        self.splitter = SentenceSplitter(chunk_size=256,chunk_overlap=25)
        self.splitter2 = SentenceSplitter(chunk_size=512,chunk_overlap=50)
        self.splitter_code = CodeSplitter(language="python", chunk_lines=40, chunk_lines_overlap=15, max_chars=1500)        
        self.response_synthesizer = get_response_synthesizer(llm = self.llm, response_mode="tree_summarize", use_async=True)
        self.query_engine_tools = []
        self.retriever_engine_tools = []
        self.summary = []

        for idx, doc in enumerate(_docs):
            idx = str(idx)
            doc_type = doc[0].metadata['resource']
            if doc[0].metadata['title'].split('.')[-1] == 'py':
                self.create_tools(idx, doc, doc_type, self.splitter_code)
            elif doc[0].metadata['title'].split('.')[-1] == 'pdf':
                self.create_tools(idx, doc, doc_type, self.splitter2)
            elif doc[0].metadata['title'].split('.')[-1] == 'pptx':
                self.create_tools(idx, doc, doc_type, self.splitter)                
            elif doc[0].metadata['title'].split('.')[-1] == 'docx':
                self.create_tools(idx, doc, doc_type, self.splitter)
            elif doc[0].metadata['resource'] =='web_page': 
                self.create_tools(idx, doc, doc_type, self.splitter)
            elif doc[0].metadata['resource'] =='youtube':    
                self.create_tools(idx, doc, doc_type, self.splitter2)

    def __len__(self):
        return len(self._docs)

    def create_tools(self, idx, doc, doc_type, splitter):
        self.index_node = GPTVectorStoreIndex.from_documents(doc, service_context = self.service_context, transformations= [splitter],show_progress=True)            
        retriever_engine = self.index_node.as_retriever()
        query_engine = self.index_node.as_query_engine()
        self.retriever_engine_tools.append(RetrieverTool.from_defaults(retriever=retriever_engine, name= doc_type + "_" + idx,
                                                                    description="Please answer questions about the {} content of the {}".format(doc_type, doc[0].metadata['title'])))

        self.query_engine_tools.append(QueryEngineTool(query_engine=query_engine, metadata=ToolMetadata(name="youtube_" + idx,
                                                                    description="Please answer questions about the {} content of the {}".format(doc_type, doc[0].metadata['title']))))     
        self.summary.append(doc[0].metadata['title'] + '\n\n' + mrs.create_document_summary('\n'.join([i.text for i in doc])) + '\n\n\n')            
    
    def retriever_documents(self):
        cohere_rerank = CohereRerank(api_key='Gq7CC5ShzvSVm1FvPrSRpdYfWt6BHfbVwCHvzRDC', top_n=2)
        bm25_retriever = BM25Retriever.from_defaults(index=self.index_node, similarity_top_k=3)
        vector_retriever = self.index_node.as_retriever(similarity_top_k=3)
        hybrid_retriever = HybridRetriever(vector_retriever, bm25_retriever)
        response_synthesizer = get_response_synthesizer(llm = self.llm,)
        custom_query_engine = RetrieverQueryEngine(
            retriever=hybrid_retriever,
            response_synthesizer=response_synthesizer,
            node_postprocessors=[cohere_rerank],
        )
        return custom_query_engine

    def hybrid_retriever_documents(self, query:str):
        """
        Answer a common question about extracted documents.
        Return retrieved texts.

        Args:
            query (str): question about extracted documents.
        """        
        cohere_rerank = CohereRerank(api_key='Gq7CC5ShzvSVm1FvPrSRpdYfWt6BHfbVwCHvzRDC', top_n=2)
        bm25_retriever = BM25Retriever.from_defaults(index=self.index_node, similarity_top_k=3)
        vector_retriever = self.index_node.as_retriever(similarity_top_k=3)
        hybrid_retriever = HybridRetriever(vector_retriever, bm25_retriever)
        response_synthesizer = get_response_synthesizer(llm = self.llm,)
        custom_query_engine = RetrieverQueryEngine(
            retriever=hybrid_retriever,
            response_synthesizer=response_synthesizer,
            node_postprocessors=[cohere_rerank],
        )
        response = custom_query_engine.query(query)
        return response.response

    def multi_index_retriever_documents(self, query:str):
        """
        Answer a common question about extracted documents.
        Return retrieved texts.

        Args:
            query (str): question about extracted documents.
        """
        query = translator(self.llm, query)
        sub_query_agent = SubQuestionQueryEngine.from_defaults(
            query_engine_tools= self.query_engine_tools,
            service_context= self.service_context,
            use_async=True,
        )                   
        res = sub_query_agent.query(query)
        return res.response

    def translate_documents(self, lang_to_translate):
        """
        Translate documents.
        Return translated texts. 

        Args:
            lang_to_translate (str): target language to translate. choose [en], [ko]. 
        """
        print(lang_to_translate)
        return self.docs
    
    def multi_retriever(self):
        cohere_rerank = CohereRerank(api_key='Gq7CC5ShzvSVm1FvPrSRpdYfWt6BHfbVwCHvzRDC', top_n=5)
        all_tools = self.query_engine_tools
        tool_mapping = SimpleToolNodeMapping.from_objects(all_tools)
        obj_index = ObjectIndex.from_objects(
            all_tools,
            tool_mapping,
            VectorStoreIndex,
            service_context = self.service_context
        )
        vector_node_retriever = obj_index.as_node_retriever(similarity_top_k=10)
        custom_node_retriever = CustomRetriever(vector_node_retriever, cohere_rerank)
        custom_obj_retriever = CustomObjectRetriever(custom_node_retriever, tool_mapping, all_tools, llm=self.llm)        
        top_agent = FnRetrieverOpenAIAgent.from_retriever(
            custom_obj_retriever,
            system_prompt=""" \
        You are an agent designed to answer queries about the documentation.
        Please always use the tools provided to answer a question. Do not rely on prior knowledge.\
        """,
            llm=self.llm,
            verbose=True,
        )      
        return top_agent