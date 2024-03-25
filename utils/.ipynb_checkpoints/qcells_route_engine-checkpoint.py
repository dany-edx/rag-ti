import chromadb
from llama_index.core.vector_stores import MetadataFilters, MetadataFilter,FilterOperator
from llama_index.core import SummaryIndex, get_response_synthesizer,  VectorStoreIndex, ServiceContext, GPTVectorStoreIndex, Document
from llama_index.vector_stores.chroma import ChromaVectorStore
from pydantic import Field, BaseModel
from llama_index.program.openai import OpenAIPydanticProgram
from typing import List, Optional
import requests
from llama_index.core.tools.tool_spec.base import BaseToolSpec
from langchain_community.document_loaders import PyPDFLoader
from llama_index.core.node_parser import SentenceSplitter
from chromadb.config import Settings
from llama_index.core.tools import QueryEngineTool, ToolMetadata
from llama_index.core.query_engine import SubQuestionQueryEngine
from llama_index.core.response_synthesizers import TreeSummarize, CompactAndRefine, Refine, Accumulate
from llama_index.core.memory import ChatMemoryBuffer
from llama_index.core.schema import NodeWithScore, Node
import json
from llama_index.core.llms import ChatMessage, MessageRole
from llama_index.postprocessor.cohere_rerank import CohereRerank
from web_catch import * 
import sys
sys.path.append('../utils')
from qcells_web_instance_search import instance_search_expanding
from yahoo_finance import * 
from llama_index.core.agent import ReActAgent
from langchain_community.document_loaders import AsyncChromiumLoader
from langchain_community.document_transformers import BeautifulSoupTransformer
import chromedriver_autoinstaller
chromedriver_autoinstaller.install()
import re
from io import BytesIO
from pptx import Presentation
from bs4 import BeautifulSoup
from llama_index.retrievers.bm25 import BM25Retriever
from llama_index.core.retrievers import BaseRetriever
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
import time 
import praw
from datetime import datetime
from llama_index.agent.openai import OpenAIAgent
from qcells_custom_rag import get_youtube_metadata, get_timeline
from youtube_transcript_api import YouTubeTranscriptApi
from llama_index.core import PromptTemplate
import chromedriver_autoinstaller
chromedriver_autoinstaller.install()
from docx import Document as py_Document
from llama_index.core.retrievers import BaseRetriever
from llama_index.agent.openai_legacy import FnRetrieverOpenAIAgent
from llama_index.core.objects import ObjectIndex,SimpleToolNodeMapping,ObjectRetriever
import multiprocessing


class global_obj(object):
    chrome_options = Options()
    chrome_options.add_argument("start-maximized")
    chrome_options.add_argument("--headless")
    chrome_options.add_argument("--no-sandbox")
    chrome_options.add_argument("--disable-dev-shm-usage")
    chrome_options.add_argument("--disable-blink-features=AutomationControlled")
    chrome_options.add_experimental_option("excludeSwitches", ["enable-automation"])
    chrome_options.add_experimental_option("detach", True)
    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36'}

class Decide_to_search_external_web(BaseModel):
    """
        Verifying if the response to the question is accurate, and deciding whether to conduct an external search.
        return
            Succeed_answer (bool) : check if answer is acceptable or mentioned real-time data.
            Decide_web_search (bool): check if need to Assess internet to answer.
            Searchable_query (str): searchable keywords on Google.(limited 4-words)
            Reason (bool): if query is clear question.
    """
    Succeed_answer: bool
    Decide_web_search: bool
    Searchable_query: str
    Reason: bool


class Result(BaseModel):
    """
    refine answer from query.
    return
        Succeed (bool): check weather answer is queried
        Url (str): which resource used
        Ingisht (str): refined answer and insight
    """
    Succeed: bool
    Url: Optional [str]
    Ingisht: str

class GoogleResult(BaseModel):
    """
    sort every information and explain it based on query.
    return
        Ingisht (str): sorted answer and explanation
    """
    Ingisht: str

def remove_duplicate_spaces(input_text):
    pattern = r'\s+'    
    result_text = re.sub(pattern, ' ', input_text)
    return result_text.strip()  # 앞뒤 공백 제거

def remove_duplicate_newlines(input_text):
    pattern = r'\n+'
    result_text = re.sub(pattern, '\n', input_text)
    return result_text.strip()  # 앞뒤 공백 제거

def translator(llm, other_lang): 
    template = (
        "You are an AI agent to translate other language into English."
        "Translate into English, if the given text is not written in English."
        "Return just tranlated word"
        "\n---------------------\n"
        "the given text: {query}"
    )
    qa_template = PromptTemplate(template)
    prompt = qa_template.format(query = other_lang)
    prompt = [ChatMessage(role='user', content=prompt)]
    translated = llm.chat(prompt)
    return translated.message.content

class HybridRetriever(BaseRetriever):
    def __init__(self, vector_retriever, bm25_retriever):
        self.vector_retriever = vector_retriever
        self.bm25_retriever = bm25_retriever
        super().__init__()

    def _retrieve(self, query, **kwargs):
        bm25_nodes = self.bm25_retriever.retrieve(query, **kwargs)
        vector_nodes = self.vector_retriever.retrieve(query, **kwargs)

        all_nodes = []
        node_ids = set()
        for n in bm25_nodes + vector_nodes:
            if n.node.node_id not in node_ids:
                all_nodes.append(n)
                node_ids.add(n.node.node_id)
        return all_nodes

class DocumentDrillDownAnalyzeToolSpec(BaseToolSpec):
    """Drill down document to query about more detail information of specific document and summarize the contents of various documents such as patents, papers, and journals."""
    spec_functions = ["document_analyzer"]
    url = '' 
    def document_analyzer(self, url:str, query:str):
        """
        Answer a query about specific documents such as patents, journals, papers, youtube etc.
        Drill down document to query about more detail information of specific document.
        Acquire data using URL links, and then query the information.

        Return resource url, output as json type. 
        Args:
            url (str): extract url from the conversation.
            query (str): question about the chosen material. if query is null, query should be "summarize"
        """    
        if query == '':
            query = 'summarize. limited 1000-words.'
        # r=requests.get(url)
        # soup = BeautifulSoup(r.content, 'html.parser')
        
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        time.sleep(3)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()

        if '.pdf' in url:
            print('PDF URL Parsing')
            if self.url != url:
                loader = PyPDFLoader(url)
                documents = loader.load_and_split()
                text_list = ' '.join([i.page_content for i in documents])
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]
        elif '.doc' in url.lower():
            print('Docx URL Parsing')
            if self.url != url:
                response = requests.get(url)
                docx_bytes = BytesIO(response.content)
                doc = py_Document(docx_bytes)
                text_list = []
                for paragraph in doc.paragraphs:
                    text_list.append(paragraph.text)
                text_list = '\n'.join(text_list)
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]  
        elif '.pptx' in url:
            if self.url != url:
                response = requests.get(url)
                pptx_data = BytesIO(response.content)
                presentation = Presentation(pptx_data)
                text_list = ""
                for i, slide in enumerate(presentation.slides):
                    text_list += f"\n\nSlide #{i}: \n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_list += f"{shape.text}\n" 
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]
        elif 'https://patents.justia.com' in url:
            if self.url != url:
                text_list = []
                for a_tag in soup.find_all('p'):
                    text_list.append(a_tag.text)
                text_list = ''.join(text_list)
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]
        elif 'https://www.cpuc.ca.gov' in url:
            try:
                if self.url != url:
                    text_list = []
                    for a_tag in soup.find_all('p'):
                        text_list.append(a_tag.text)
                    text_list = ''.join(text_list)
                    self.hybrid_retriever = self.documentize(text_list, url)
                result = self.hybrid_retriever.retrieve(query)
                result = [i.text for i in result]
            except:
                pass
        elif 'https://paperswithcode.com' in url:
            href_list = []
            web_list = []
            git_list = []
            div_soup = soup.find('div', class_='paper-abstract')            
            for a_tag in div_soup.find_all('p'):
                href_list.append(a_tag.text.strip())        
            for a_tag in soup.find_all('a'):
                if 'href' in a_tag.attrs:
                    if '.pdf' in a_tag['href']:
                        web_list.append(a_tag['href'])
                    if 'github' in a_tag['href']:
                        git_list.append(a_tag['href'])
            result = {
                'paper_summary' : href_list[0],
                'pdf_urls' : list(set(web_list)),
                'github_urls' : list(set(git_list))[:5]}
        elif 'https://www.nature.com' in url:
            if self.url != url:
                text_list = []
                for a_tag in soup.find_all('p'):
                    text_list.append(a_tag.text)
                text_list = ''.join(text_list)
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]
        elif 'https://patents.google.com' in url:
            if self.url != url:
                text_list = []
                for div_tag in soup.find_all('div'):
                    if 'num' in div_tag.attrs:
                        text_list.append(div_tag.text)
                text_list = '\n'.join(text_list)
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]
        elif 'www.reddit.com' in url:
            if self.url != url:
                reddit = praw.Reddit(client_id='GIbgvR6hff64rcsIHihC3g',
                                 client_secret='2yora5MRJLBsGLfGxj_5ooDRCf68Kw',
                                 user_agent='my_apps/1.0')
                post_id = url.split('comments/')[-1].split('/')[0]
                submission = reddit.submission(id=post_id)
                content = []
                comments =  []
                total_comments = []
                for i, comment in enumerate(submission.comments.list()):    
                    try:
                        comments.append("comment{} : {}".format(i, comment.body))    
                    except:
                        pass
                content.append([submission.title, submission.url,submission.selftext, comments[:3]])
                result = [{'title': item[0], 'url': item[1], 'content': item[2], 'total_comments': item[3]} for item in content]
        elif 'www.youtube.com' in url:
            if self.url != url:
                meta_data = get_youtube_metadata(url)
                data = YouTubeTranscriptApi.get_transcript(url.split('v=')[-1])
                documents = []
                for i in data:
                    i['div'] = int(i['start'] / 60)  
                text_list= []
                distinct_div = set(item['div'] for item in data)
                for idx, d in enumerate(distinct_div):
                    texts = []
                    k = [i for i in data if i['div'] == d]
                    start_time_min = min([i['start'] for i in k])
                    start_time_max = max([i['start'] for i in k])
                    text_div = [i['text'] for i in k]
                    texts.append('[youtube play time] {}\n'.format(get_timeline(start_time_min)) + ' '.join(text_div))
                    texts = '\n'.join(texts)
                    text_list.append(texts)            
                text_list = '\n'.join(text_list)
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]
        else:
            if self.url != url:
                text_list = []
                for a_tag in soup.find_all('p'):
                    text_list.append(a_tag.text)
                text_list = ''.join(text_list)
                self.hybrid_retriever = self.documentize(text_list, url)
            result = self.hybrid_retriever.retrieve(query)
            result = [i.text for i in result]

        self.url = url
        
        try:
            return {"resource_url": self.url, "ouput":result}
        except:
            return {"resource_url": "url not found", "output":'Retry. please extract proper url again.'}
    
    def documentize(self, text_list, url):
        patent_data_documents = [Document(text=text_list, metadata = {'url' : url})]
        splitter = SentenceSplitter(chunk_size=256,chunk_overlap=20)
        patent_index = VectorStoreIndex.from_documents(documents=patent_data_documents, transformations=[splitter], service_context = self.service_context, show_progress=True)
        bm25_retriever = BM25Retriever.from_defaults(index=patent_index, similarity_top_k=1)
        vector_retriever = patent_index.as_retriever(similarity_top_k=1)
        hybrid_retriever = HybridRetriever(vector_retriever, bm25_retriever)
        return hybrid_retriever

class CaliforniaUtilityCommisionSearchToolSepc(DocumentDrillDownAnalyzeToolSpec, BaseToolSpec):
    # spec_functions = ['commission_search']
    """California public solar energy utility commission forum"""
    def california_utility_commission_search(self, query:str):
        """
        Search California energy commision forum for policies, commissions and regulations.
        Return title, abstract and url as json.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        url = 'https://www.cpuc.ca.gov/search#q={}&sort=relevancy'.format(query)
        driver = webdriver.Chrome( options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        time.sleep(3)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.close()
        
        div_tags = soup.find_all('div', {'class':'coveo-list-layout CoveoResult'})
        text_list = []
        for div in div_tags:
            a_tags = div.find_all('a')
            text_list.append([a_tags[0].text, div.text, a_tags[0].attrs['href']])
        result_dict_list = [{'title': item[0], 'abstract': item[1], 'url': item[2]} for item in text_list]
        return result_dict_list[:5]

class CaliforniaEnergyMarketSearchToolSpec(BaseToolSpec):
    """California States Energy grid market, policy, trend and regulation search tool spec."""
    def california_information_search(self, query):
        """
        Answer about real-time power data including generation, consumption, and prices, market reports for trend analysis, operational notices for grid management, 
        regulations, resources, and tools for market participants, and information on events and education programs.

        Return title, url and abstract as json.
        Args:
            query (str): searchable words on google (limited 4 words)
        """    
        query = translator(self.llm, query)
        url= 'https://www.caiso.com/search/Pages/WWW-Not-Archived-Results.aspx?k={}'.format(query.replace(' ','+'))
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        time.sleep(5)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()
        
        title = []
        urls = []
        abstract = []
        main_div = soup.find_all('div', {'class':'ms-srch-group-content'})
        div_tags = main_div[0].find_all('div', {'name' : 'Item'})
        for i in div_tags:
            a_tag = i.find_all('a')[0]
            urls.append(a_tag.attrs['href'].split('#')[0])
            title.append(a_tag.text)
            abstract.append(i.find_all('div', {'ms-srch-item-summary'})[0].text)
        
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]
        return result_dict_list[:5]
        
class TexasEnergyMarketSearchToolSpec(BaseToolSpec):
    """Texas States Energy grid market, policy, trend and regulation search tool spec."""
    def texas_information_search(self, query):
        """
        Answer about real-time power data including generation, consumption, and prices, market reports for trend analysis, operational notices for grid management, 
        regulations, resources, and tools for market participants, and information on events and education programs.

        Return title, url and abstract as json.
        Args:
            query (str): searchable words on google (limited 4 words)
        """    
        query = translator(self.llm, query)
        url= 'https://www.ercot.com/search?q={}'.format(query.replace(' ','+'))
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        time.sleep(5)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()
        
        title = []
        urls = []
        abstract = []
        div_tags = soup.find_all('div', id = 'search-results')
        for a in div_tags[0].find_all(['div','a']):
            if 'href' in a.attrs:
                if 'www.ercot.com' in a.attrs['href']:
                    urls.append(a.attrs['href'])    
                else:
                    urls.append('https://www.ercot.com' + a.attrs['href'])
                title.append(a.text)
            if 'class' in a.attrs:
                if 'my-2' in a['class']:
                    abstract.append(a.text)
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]
        return result_dict_list[:5]

class TexasUtilityCommissionSearchToolSpec(BaseToolSpec):
    """Texas public solar energy utility commission forum"""
    def texas_utility_commission_search(self, query):
        """
        Answer about real-time power data including generation, consumption, and prices, market reports for trend analysis, operational notices for grid management, 
        regulations, resources, and tools for market participants, and information on events and education programs.

        Return title, url and abstract as json.
        Args:
            query (str): searchable words when search dissertation (limited 4 words)
        """
        query = translator(self.llm, query)
        url= 'https://www.puc.texas.gov/agency/sitesearch.aspx?q={}'.format(query.replace(' ','+'))
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()
        title = []
        urls = []
        abstract = []        
        divs = soup.find_all('div', {'class':'gsc-webResult gsc-result'})
        for d in divs:
            a_tags = d.find_all('a')[0]
            title.append(a_tags.text)
            urls.append(a_tags.attrs['href'])
            d_tag = d.find_all('div', {'class' : 'gs-bidi-start-align gs-snippet'})[0]
            abstract.append(d_tag.text)
        
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]
        return result_dict_list[:5]
        

class RedditMarketingSearchToolSpec(BaseToolSpec):
    """SNS Reddit marketing search tool spec."""
    def reddit_post_search(self, query:str):
        """
        Search SNS Reddit about solar marketing information. answer a query about solar panel users opinions.
        Return post title, url, content and comments as json.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        reddit = praw.Reddit(client_id='GIbgvR6hff64rcsIHihC3g',
                             client_secret='2yora5MRJLBsGLfGxj_5ooDRCf68Kw',
                             user_agent='my_apps/1.0')
        subreddit = reddit.subreddit('solar')
        search_results = subreddit.search(query, limit=3, time_filter= 'month')
        content = []
        total_comments = []
        for idx, submission in enumerate(search_results):
            comments = []
            for i, comment in enumerate(submission.comments.list()):    
                try:
                    comments.append("comment{} : {}".format(i, comment.body))    
                except:
                    pass
            url = 'https://www.reddit.com/r/solar/comments/' + submission.id
            content.append([submission.title, url, submission.selftext[:2000], comments[:2]])
        result_dict_list = [{'title': item[0], 'url': item[1], 'content': item[2], 'total_comments': item[3]} for item in content]
        return result_dict_list[:4]

class AzureCloudManualSearchToolSpec(BaseToolSpec):
    """Azure microsoft cloud serivce manual search tool spec."""
    def azure_cloud_manaul(self, query:str):
        """
        Search Azure manaul search. 
        Return manual title, url, abstract as json.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        url= 'https://learn.microsoft.com/en-us/search/?terms={}'.format(query.replace(' ','+'))
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url)
        time.sleep(3)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()
        
        title = []
        urls = []
        abstract = []
        
        div_tags = soup.find_all('li', {'class':'margin-top-sm'})
        for d in div_tags:
            a_tag = d.find_all('a')
            p_tag = d.find_all('p')
        
            urls.append(a_tag[0].attrs['href'])
            title.append(a_tag[0].text.strip())
            abstract.append(p_tag[0].text)
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]
        return result_dict_list[:5]

class IBMDevSearchToolSpec(BaseToolSpec):
    """IBM developer manual search tool spec."""
    def ibm_dev_manaul(self, query:str):
        """
        Search IBM manaul search. 
        Return manual title, url, abstract as json.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        url= 'https://developer.ibm.com/?q={}'.format(query.replace(' ','+'))
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        time.sleep(5)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()
        
        title = []
        urls = []
        abstract = []
        div_tags = soup.find_all('div', {'class':'ibm--result-item'})
        for d in div_tags:
            a_tag = d.find_all('a')[0]
            t = d.find_all('h2', {'class' : 'ibm--result-item__title'})[0].text
            abs = d.find_all('p', {'class' : 'ibm--result-item__description'})[0].text
            urls.append(a_tag.attrs['href'])
            title.append(t)
            abstract.append(abs)
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]    
        return result_dict_list[:5]
        
class AWSCloudManualSearchToolSpec(BaseToolSpec):
    """AWS cloud serivce manual search tool spec."""
    def aws_cloud_manaul(self, query:str):
        """
        Search AWS manaul search. 
        Return manual title, url, abstract as json.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        url= 'https://aws.amazon.com/ko/search/?searchQuery={}#facet_type=documentation&page=1'.format(query.replace(' ','+'))
        driver = webdriver.Chrome(options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url)
        time.sleep(3)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        driver.quit()
        div_tags = soup.find_all('li', {'class':'item documentation clearfix'})
        
        title = []
        urls = []
        abstract = []
        
        for d in div_tags:
            title.append(d.find_all('a')[0].text)
            urls.append(d.find_all('a')[0].attrs['href'])
            abstract.append(d.find_all('div', {'class':'lb-item-desc'})[0].text)
        
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]
        return result_dict_list[:5]
        
class ComputerSciencePaperSearchToolSpec(DocumentDrillDownAnalyzeToolSpec, BaseToolSpec):
    """computer science related paper search tool spec."""
    def computer_science_paper_search(self, query:str, is_more=False):
        """
        Search IT and AI tech related paper. answer paper url parsed link and name.
        Return Urls and paper titles.
        
        Args:
            query (str): searchable words when search dissertation (limited 4 words)
            is_more [Optional] (bool): Count up when the user fetch more data about the same topic or query.  
        """
        query = translator(self.llm, query)
        url = 'https://paperswithcode.com/search?q_meta=&q_type=&q={}'.format(query.replace(' ', '+'))
        self.driver = webdriver.Chrome(options=global_obj.chrome_options)
        self.driver.get(url) 
        # if is_more == False:
        #     self.prev_len = 0
        #     self.is_more_cnt = 0
        # else:
        #     self.is_more_cnt = self.is_more_cnt + 1

        # for i in range(self.is_more_cnt):
        #     if self.prev_len > 0:
        #         new_height = self.driver.execute_script('return document.body.scrollHeight')
        #         self.driver.execute_script("window.scrollTo(0, {})".format(new_height))
        #         time.sleep(5)
        time.sleep(2)
        soup = BeautifulSoup(self.driver.page_source, 'html.parser')
        title = []
        urls = []
        abstract = []
        div_tags = soup.find_all('div', {'class':'col-lg-9 item-content'})
        for d in div_tags:
            a_tag = d.find_all('a')
            p_tag = d.find_all('p', {'class' : 'item-strip-abstract'})
            for a,p in zip(a_tag, p_tag):
                if 'href' in a.attrs:
                    if '/paper/' in a.attrs['href']:
                        if '#' not in a.attrs['href']:
                            urls.append('https://paperswithcode.com' + a.attrs['href'])
                            title.append(a.text)
                abstract.append(p.text)
        content = []
        for t, u, a in zip(title, urls, abstract):
            content.append([t, u, a])
        result_dict_list = [{'title': item[0], 'url': item[1], 'abstract': item[2]} for item in content]                
        # href_list = href_list[self.prev_len:]
        # self.prev_len = self.prev_len + len(href_list)        

        # href_list = self.removal_duplicates(href_list)
        # result_dict_list = [{'paper_url': item[0], 'paper_name': item[1]} for item in href_list]                                
        self.driver.quit()

        return result_dict_list[:5]
        
    def removal_duplicates(self, href_list):
        grouped_data = {}
        for key, value in href_list:
            if key not in grouped_data:
                grouped_data[key] = []
            grouped_data[key].append(value)
        
        filtered_data = []
        for key, values in grouped_data.items():
            max_value = max(values, key=len)
            filtered_data.append([key, max_value])
        return filtered_data
        
class ChemistryEngineeringJournalSearchToolSpec(DocumentDrillDownAnalyzeToolSpec, BaseToolSpec):
    """Journal random search tool spec."""
    def scholar_paper_search(self, query):
        """
        Search for general topic's academical paper names and ids.
        Return paper titles and urls as json type.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        url= 'https://scholar.google.com/scholar?as_sdt=2007&q={}&hl=en&as_ylo=2024'.format(query.replace(' ','+'))
        r=requests.get(url)
        soup = BeautifulSoup(r.content, 'html.parser')
        anchor_elements = soup.find_all('a')
        href_list = []
        for a_tag in soup.find_all('a'):
            if 'href' in a_tag.attrs:
                if 'http' in a_tag['href']:
                    if 'google' not in a_tag['href']:
                        span_tags = a_tag.find_all('span')
                        if not span_tags:
                            href_list.append([a_tag['href'], a_tag.text])
    
        result_list_of_dicts = [{'url': item[0].replace('Link: ', ''), 'Title': item[1]} for item in href_list]
        return result_list_of_dicts[:5]
    
    def chemistry_journal_search(self, query):
        """
        Searching journal names and ids related to chemical engineering.
        Return journal names and url as json type.
        
        Args:
            query (str): searchable words on google(limited 4 words).
        """
        query = translator(self.llm, query)
        url= 'https://www.nature.com/search?q={}&date_range=last_year&order=relevance'.format(query.replace(' ','+'))
        r=requests.get(url)
        soup = BeautifulSoup(r.content, 'html.parser')
        href_list = []
        for a_tag in soup.find_all('a'):
            if 'href' in a_tag.attrs:
                if 'article' in a_tag['href']:
                    href_list.append(['https://www.nature.com' + a_tag['href'], a_tag.text])
        
        result_list_of_dicts = [{'url': item[0].replace('Link: ', ''), 'Title': item[1]} for item in href_list]
        return result_list_of_dicts[:5]
    
class JustiaPatentRandomSearchToolSpec(DocumentDrillDownAnalyzeToolSpec, BaseToolSpec):
    """Web Patent random search tool spec using Jistia."""
    def justia_patent_search(self, query):
        """
        Search patents on justia web for specific topic's patent url, abstract and title.
        Please if there no result, go to Google patent search. 
        
        Return patent url, abstract, title as json type.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        query = translator(self.llm, query)
        url= 'https://patents.justia.com/search?q={}'.format(query.replace(' ','+'))
        r=requests.get(url)
        soup = BeautifulSoup(r.content, 'html.parser')
        
        href_list = []
        title_list = []
        abstract_list = []
        
        for i, div_tag in enumerate(soup.find_all('div', attrs = {'class':['head','abstract']})):
            if 'abstract' in div_tag.attrs['class']:
                abstract_list.append(div_tag.text)
            
            for a_tag in div_tag.find_all('a'):
                if 'href' in a_tag.attrs:
                    if '/patent/' in a_tag['href']:
                        href_list.append('https://patents.justia.com' + a_tag['href'])
                        title_list.append(a_tag.text)
        merge_list = list(zip(href_list, title_list, abstract_list))
        result_list_of_dicts = [{'Url': item[0].replace('Link: ', ''), 'Title': item[1], 'Abstract': item[2]} for item in merge_list]                    
        return result_list_of_dicts[:3]

class GooglePatentRandomSearchToolSpec(DocumentDrillDownAnalyzeToolSpec, BaseToolSpec):
    """Patents search second tool spec using google-patent."""
    def google_patent_search(self, query:str):    
        """
        Search patents on google.patent web for specific topic's patent names and ids.
        Return patent pdf, url link, title as json type.
    
        Args:
            query (str): searchable words on google (limited 4 words)
        """    
        query = translator(self.llm, query)
        url = 'https://patents.google.com/?q=({})'.format(query.replace(' ','+'))
        driver = webdriver.Chrome( options=global_obj.chrome_options)
        driver.delete_all_cookies()
        driver.get(url) 
        time.sleep(3)
        soup = BeautifulSoup(driver.page_source, 'html.parser')
        href_list = []
        for item_tag in soup.find_all('search-result-item'):
            a_tag = item_tag.find_all('a')
            for a in a_tag:
                if 'href' in a.attrs:
                    if 'https:' in a['href']:
                        url = a['href']
        
            span_tag = item_tag.find_all('span')
            for s in span_tag:
                if 'data-proto' in s.attrs:
                    if s['data-proto'] == 'OPEN_PATENT_PDF':
                        id_name = 'https://patents.google.com/patent/' + s.text
            texts = []
            for s in span_tag:
                if 'id' in s.attrs:
                    if s['id'] == 'htmlContent':
                        texts.append(s.text)
            texts = ''.join(texts)
            href_list.append([url, id_name, texts])
        result_dict_list = [{'pdf': item[0], 'url': item[1], 'title': item[2]} for item in href_list]                            
        driver.delete_all_cookies()
        driver.quit()
        return result_dict_list[:3]

class GoogleDeepSearchToolSpec(BaseToolSpec):
    """Google search for high level questions."""
    spec_functions = ["get_answer"]
    result_queue = multiprocessing.Queue()
    index = ''
    
    def get_answer(self, query):
        """
        Answer a detail analyzed query. 
        Return analyzed detail retrieved data.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        self.query = query
        self.temperal_index()
        
        bm25_retriever = BM25Retriever.from_defaults(index=self.index, similarity_top_k=1)
        vector_retriever = self.index.as_retriever(similarity_top_k=1)
        hybrid_retriever = HybridRetriever(vector_retriever, bm25_retriever)
        result = hybrid_retriever.retrieve(query)
        return [i.text for i in result]
    
    def get_soup(self):
        url= 'https://www.google.com/search?q={}&tbs=qdr:12m'.format(self.query.replace(' ','+'))
        url_news = 'https://www.google.com/search?sca_esv=%EB%89%B4%EC%8A%A4&q={}&tbm=nws'.format(self.query.replace(' ','+'))
        
        try:
            driver = webdriver.Chrome(options=global_obj.chrome_options)
            driver.delete_all_cookies()
            driver.set_page_load_timeout(10)

            driver.get(url)
            time.sleep(1)
            self.soup = BeautifulSoup(driver.page_source, 'html.parser')
            
            driver.get(url_news)
            time.sleep(1)
            self.soup_news = BeautifulSoup(driver.page_source, 'html.parser')
            driver.quit()
        except:
            driver.quit()
            print('failed pass!')
            pass
        return self.soup, self.soup_news

    def google_search_q_instance(self, soup):
        urls = []
        for div in soup.find_all('div'):
            if 'data-initq' in div.attrs:
                a_tags = div.find_all('a')
                for a in a_tags:
                    if 'http' in a.attrs['href']:
                        urls.append(a.attrs['href'])
        content = []
        t1 = time.time()
        processes = []
        for url in urls[:2]:
            p = multiprocessing.Process(target=self.get_string_from_news, kwargs={'url': url})
            processes.append(p)
            p.start()
        for p in processes:
            p.join()
        p.close()
        results = []
        while not self.result_queue.empty():
            result_string = self.result_queue.get()
            results.append(result_string)        
        return results
        
    def google_search_instance(self, soup):
        print('google_search_instance...')
        anchor_elements = [] 
        for i in soup.find_all('div', id = 'search')[0].find_all(['div']):
            if len(i.text) >2:
                if len(i.attrs) == 1:
                    if 'jsslot' not in i.attrs:
                        if 'class' in i.attrs:
                            if len(i.attrs['class']) == 1:
                                if len(i.find_all('script')) == 0:
                                    anchor_elements.append(i.text)
        
        a = list(set([i for i in anchor_elements if len(i) > 2]))
        sorted_text_list = sorted(a, key=len)
        result_text = remove_duplicate_spaces(''.join(sorted_text_list[-3:]))
        result_text = remove_duplicate_newlines(result_text)
        program = OpenAIPydanticProgram.from_defaults(
            output_cls=GoogleResult,
            llm=self.llm,
            prompt_template_str=(
                "This is google search result context.\n"
                "I will give you question and google reult\n"
        
                "This is question\n"
                "{query}\n"
                
                "\n"        
                "This is Google result\n"
                "{context}"
            ),
        )        
        res = program(query = self.query, context = result_text)      
        return [Document(text=res.Ingisht, metadata = {'url' : 'instance_search'})]
    
    def get_string_from_news(self, url):
        t1 = time.time()
        try:
            driver = webdriver.Chrome(options=global_obj.chrome_options)
            driver.delete_all_cookies()
            driver.set_page_load_timeout(5)
            driver.get(url)
            
            soup = BeautifulSoup(driver.page_source, 'html.parser')
            anchor_elements = soup.find_all('p')
            result_string = '\n'.join([i.text.strip() for i in anchor_elements])  
            driver.quit()
            self.result_queue.put(Document(text=result_string[:20000], metadata = {'url' : url}))
            print(url, time.time() - t1)
        except:
            driver.quit()
            pass
        
    def google_search_news_instance(self, soup):
        title_list = []
        href_list = []
        contents = []
        d_tag = soup.find_all('div', id = 'search')
        d_tag = d_tag[0].find_all('div')[0]
        for d in d_tag:
            a_tag = d.find_all('a')
            if len(a_tag) > 0:
                for a in a_tag:
                    if 'href' in a.attrs:
                        href_list.append(a.attrs['href'])
                        title_list.append(a.text)
        content = []
        for t, u in zip(title_list, href_list):
            content.append([t, u])
        result_dict_list = [{'abstract': item[0], 'url': item[1]} for item in content]                
        t1 = time.time()
        urls = [c['url'] for c in result_dict_list]    
        processes = []
        for url in urls[:5]:
            p = multiprocessing.Process(target=self.get_string_from_news, kwargs={'url': url})
            processes.append(p)
            p.start()
        for p in processes:
            p.join()
        p.close()
        
        results = []
        while not self.result_queue.empty():
            result_string = self.result_queue.get()
            results.append(result_string)
        
        print(time.time() - t1)
        return results
        
    def temperal_index(self):
        splitter = SentenceSplitter(chunk_size=512,chunk_overlap=51)
        self.soup, self.soup_news = self.get_soup()
        documents = self.google_search_instance(self.soup) + self.google_search_news_instance(self.soup_news)  
        self.index = VectorStoreIndex.from_documents(documents=documents, transformations=[splitter], service_context = self.service_context, show_progress=True)


class GoogleRandomSearchToolSpec(ChemistryEngineeringJournalSearchToolSpec, ComputerSciencePaperSearchToolSpec, JustiaPatentRandomSearchToolSpec, GooglePatentRandomSearchToolSpec, RedditMarketingSearchToolSpec, TexasEnergyMarketSearchToolSpec, CaliforniaUtilityCommisionSearchToolSepc, TexasUtilityCommissionSearchToolSpec, AWSCloudManualSearchToolSpec, AzureCloudManualSearchToolSpec, CaliforniaEnergyMarketSearchToolSpec, IBMDevSearchToolSpec):
    """Google general information random search tool spec."""
    spec_functions = ["get_instinct_search_url"]
    def get_instinct_search_url(self, query):
        """
        Answer a simple question using google search result. 
        Return simple answer.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        url= 'https://www.google.com/search?q={}&tbs=qdr:6m'.format(query.replace(' ','+'))
        r=requests.get(url)
        soup = BeautifulSoup(r.content, 'html.parser')
        anchor_elements = soup.find_all('div')
        a = list(set([i.text for i in anchor_elements if len(i.text) > 2]))
        sorted_text_list = sorted(a, key=len)
        result_text = remove_duplicate_spaces(sorted_text_list[-1])
        result_text = remove_duplicate_newlines(result_text)
        program = OpenAIPydanticProgram.from_defaults(
            output_cls=GoogleResult,
            llm=self.llm,
            prompt_template_str=(
                "This is google search result context.\n"
                "I will give you question and google reult\n"
        
                "This is question\n"
                "{query}\n"
                
                "\n"        
                "This is Google result\n"
                "{context}"
            ),
            verbose=True,
        )        
        res = program(query = query, context = result_text)        
        return res.Ingisht
    
    def get_google_searched_urls(self, url):
        r=requests.get(url)
        soup = BeautifulSoup(r.content, 'html.parser')
        anchor_elements = soup.find_all('a')
        
        access_url = []
        for element in anchor_elements:
            href=element.get('href')
            if 'https' in href:
                if not 'google' in href:
                    if not 'youtube' in href:
                        if not 'amazon' in href:
                            if not 'yahoo' in href:
                                urls = ''.join(href.split('=')[1:])
                                urls = urls.split('&')[0]
                                if len(urls.split('/')) > 4:
                                    access_url.append(urls)
        return access_url
    
    def google_search_drill_down_analysis(self, query):
        """
        Answer a detail analyzed query. use this only when get_instinct_search_url failed to answer. 
        Return analyzed detail retrieved data.
        
        Args:
            query (str): searchable words on google (limited 4 words)
        """
        url= 'https://www.google.com/search?q={}&tbs=qdr:6m'.format(query.replace(' ','+'))
        url_news = 'https://www.google.com/search?sca_esv=%EB%89%B4%EC%8A%A4&q={}&tbs=qdr:m&tbm=nws'.format(query.replace(' ','+'))
        access_url = self.get_google_searched_urls(url)[1:4] + self.get_google_searched_urls(url_news)[1:4]
        print(access_url)
        loader = AsyncChromiumLoader(access_url)
        bs_transformer = BeautifulSoupTransformer()
        docs_transformed = bs_transformer.transform_documents(loader.load(), unwanted_tags= ['a','style','script'], tags_to_extract=['p','span'])
        web_data_documents = [Document(text=i.page_content) for i in docs_transformed if len(i.page_content) > 100]
        splitter = SentenceSplitter(chunk_size=256,chunk_overlap=20)
        web_index = GPTVectorStoreIndex.from_documents(documents=web_data_documents, service_context = self.service_context, transformations=[splitter], show_progress=True)
        bm25_retriever = BM25Retriever.from_defaults(index=web_index, similarity_top_k=1)
        vector_retriever = web_index.as_retriever(similarity_top_k=1)
        hybrid_retriever = HybridRetriever(vector_retriever, bm25_retriever)
        result = hybrid_retriever.retrieve(query)
        return [i.text for i in result]

class CustomRetriever(BaseRetriever):
    def __init__(self, vector_retriever, postprocessor=None):
        self._vector_retriever = vector_retriever
        self._postprocessor = postprocessor 
        super().__init__()

    def _retrieve(self, query_bundle):
        retrieved_nodes = self._vector_retriever.retrieve(query_bundle)
        filtered_nodes = self._postprocessor.postprocess_nodes(
            retrieved_nodes, query_bundle=query_bundle
        )
        return filtered_nodes

class CustomObjectRetriever(ObjectRetriever):
    def __init__(self, retriever, object_node_mapping, all_tools, llm=None):
        self._retriever = retriever
        self._object_node_mapping = object_node_mapping
        self._llm = llm 

    def retrieve(self, query_bundle):
        nodes = self._retriever.retrieve(query_bundle)
        tools = [self._object_node_mapping.from_node(n.node) for n in nodes]

        sub_question_engine = SubQuestionQueryEngine.from_defaults(
            query_engine_tools=tools, llm=self._llm
        )
        sub_question_description = f"""\
Useful for any queries that involve comparing multiple documents. ALWAYS use this tool for comparison queries - make sure to call this \
tool with the original query. Do NOT use the other tools for any queries involving multiple documents.
"""
        sub_question_tool = QueryEngineTool(
            query_engine=sub_question_engine,
            metadata=ToolMetadata(
                name="compare_tool", description=sub_question_description
            ),
        )
        return tools + [sub_question_tool]

class VectordbSearchToolSpec(GoogleRandomSearchToolSpec):
    spec_functions = ["vector_database_search", "google_patent_search", "justia_patent_search",  "chemistry_journal_search",  "computer_science_paper_search",  "get_instinct_search_url",  "reddit_post_search", "document_analyzer", "texas_information_search", "california_utility_commission_search", 'texas_utility_commission_search', 'aws_cloud_manaul', 'azure_cloud_manaul', "california_information_search", "ibm_dev_manaul"]
    def __init__(self, llm, service_context):
        self.llm = llm
        self.service_context = service_context
        self.response_synth = []
        self.splitter = SentenceSplitter(chunk_size=512,chunk_overlap=20)
        self.pdf_url = ''
        self.cohere_rerank = CohereRerank(api_key='Gq7CC5ShzvSVm1FvPrSRpdYfWt6BHfbVwCHvzRDC', top_n=4)
        self.filters = MetadataFilters(
            filters=[
                MetadataFilter(key="release_date_year", operator=FilterOperator.GTE, value=2024),
                MetadataFilter(key="release_date_year", operator=FilterOperator.LT, value=2025),
            ]
        )
        self.filters1 = MetadataFilters(
            filters=[
                MetadataFilter(key="release_date_year", operator=FilterOperator.GTE, value=2023),
                MetadataFilter(key="release_date_year", operator=FilterOperator.LT, value=2024),
            ]
        )
        self.filters2 = MetadataFilters(
            filters=[
                MetadataFilter(key="release_date_year", operator=FilterOperator.GTE, value=2022),
                MetadataFilter(key="release_date_year", operator=FilterOperator.LT, value=2023),
            ]
        )
        self.filters3 = MetadataFilters(
            filters=[
                MetadataFilter(key="release_date_year", operator=FilterOperator.GTE, value=2021),
                MetadataFilter(key="release_date_year", operator=FilterOperator.LT, value=2022),
            ]
        )
        self.connect_db()
        self.create_indexes()
        self.create_tools()
        
    def connect_db(self):
        db = chromadb.HttpClient(host='40.125.73.184', port=8000, settings=Settings(chroma_client_auth_provider="chromadb.auth.basic.BasicAuthClientProvider",chroma_client_auth_credentials="qcells:qcells"))
        chroma_collection = db.get_collection("pv_magazine_sentence_split")
        vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
        self.index = VectorStoreIndex.from_vector_store(vector_store,service_context=self.service_context, similarity_top_k=4, use_async=True)       

    def create_indexes(self):
        self.vector_queryengine = self.index.as_query_engine(vector_store_query_mode="hybrid", alpha = 0.5 ,node_postprocessors=[self.cohere_rerank], similarity_top_k=10, filters = self.filters, verbose = True)
        self.vector_queryengine1 = self.index.as_query_engine(vector_store_query_mode="hybrid", alpha = 0.5 ,node_postprocessors=[self.cohere_rerank], similarity_top_k=10, filters = self.filters1, verbose = True)
        self.vector_queryengine2 = self.index.as_query_engine(vector_store_query_mode="hybrid", alpha = 0.5 ,node_postprocessors=[self.cohere_rerank],similarity_top_k=10, filters = self.filters2, verbose = True)
        self.vector_queryengine3 = self.index.as_query_engine(vector_store_query_mode="hybrid", alpha = 0.5 ,node_postprocessors=[self.cohere_rerank], similarity_top_k=10, filters = self.filters3, verbose = True)

    
    def create_tools(self):
        self.query_engine_tools = [
            QueryEngineTool(
                query_engine=self.vector_queryengine,
                metadata=ToolMetadata(
                    name="news_retriever_in_2024",
                    description="Provides Renewable Energy industry related news documents that are released in 2024. this documents are stored in vector database."
                ),
            ),
            QueryEngineTool(
                query_engine=self.vector_queryengine1,
                metadata=ToolMetadata(
                    name="news_retriever_in_2023",
                    description="Provides Renewable Energy industry related news documents that are released in 2023. this documents are stored in vector database."
                ),
            ),
            QueryEngineTool(
                query_engine=self.vector_queryengine2,
                metadata=ToolMetadata(
                    name="news_retriever_in_2022",
                    description="Provides Renewable Energy industry related news documents that are released in 2022. this documents are stored in vector database."
                ),
            ),
            QueryEngineTool(
                query_engine=self.vector_queryengine3,
                metadata=ToolMetadata(
                    name="news_retriever_in_2021",
                    description="Provides Renewable Energy industry related news documents that are released in 2021. this documents are stored in vector database."
                ),
            ),
        ]         

    def vector_database_search(self, query: str):
        """
        Answer a query about renewable energy industry news articles like trends, tech, new products release, manufacturing and businesses
        
        Return query results
        Args:
            query (str): question about renewable energy industry news
        """              
        query = translator(self.llm, query)
        all_tools = self.query_engine_tools
        tool_mapping = SimpleToolNodeMapping.from_objects(all_tools)
        obj_index = ObjectIndex.from_objects(
            all_tools,
            tool_mapping,
            VectorStoreIndex,
            service_context = self.service_context
        )
        vector_node_retriever = obj_index.as_node_retriever(similarity_top_k=10)
        custom_node_retriever = CustomRetriever(vector_node_retriever, self.cohere_rerank)
        custom_obj_retriever = CustomObjectRetriever(custom_node_retriever, tool_mapping, all_tools, llm=self.llm)        
        top_agent = FnRetrieverOpenAIAgent.from_retriever(
            custom_obj_retriever,
            system_prompt=""" \
        You are an agent designed to answer queries about the solar energy industry.
        Please always use the tools provided to answer a question. Do not rely on prior knowledge.\
        """,
            llm=self.llm,
            verbose=True,
        )  
        res = top_agent.query(query)
        return res.response 


def qcell_engine(llm, embedding):
    memory = ChatMemoryBuffer.from_defaults(token_limit=4000)
    service_context = ServiceContext.from_defaults(llm=llm,embed_model=embedding)
    vector_tool_spec = VectordbSearchToolSpec(llm, service_context)
    chat_engine = ReActAgent.from_llm(vector_tool_spec.to_tool_list(), memory=memory, max_iterations = 10, llm = llm, verbose = True)
    return chat_engine

def web_engine(llm, embedding):
    service_context = ServiceContext.from_defaults(llm=llm,embed_model=embedding)
    memory = ChatMemoryBuffer.from_defaults(token_limit=4000)
    web_tool_spec = GoogleRandomSearchToolSpec()
    web_tool_spec.service_context = service_context
    web_tool_spec.llm = llm
    chat_engine = OpenAIAgent.from_llm(web_tool_spec.to_tool_list(),
                                       prefix_messages=[ChatMessage(role="system", content="Check if the user query need to use google search or not. if it need, use google search tool.")],
                                       memory=memory, max_iterations = 10, llm = llm, verbose = True)
    return chat_engine

def high_level_engine(llm, embedding):
    service_context = ServiceContext.from_defaults(llm=llm,embed_model=embedding)
    memory = ChatMemoryBuffer.from_defaults(token_limit=4000)
    da = GoogleDeepSearchToolSpec()
    da.service_context = service_context
    da.llm = llm
    rag = ReActAgent.from_llm(da.to_tool_list(), memory=memory, max_iterations = 10, llm = llm, verbose = True)
    return rag
# if __name__:
#     chat_engine = qcell_engine()